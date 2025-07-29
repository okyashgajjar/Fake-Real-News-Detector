import streamlit as st
import re
import pickle
import fitz  # Import PyMuPDF
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.linear_model import RidgeClassifier
from docx import Document
from pptx import Presentation

# ---------- Load model and vectorizer ----------
# Ensure these .pkl files are in the same directory as your app.py
# You might need to adjust the paths if they are located elsewhere.
try:
    with open("tfid_Text_vectorizer.pkl", "rb") as f:
        vectorizer = pickle.load(f)
    with open("News_Predictor_model_RidgeClassifier_.pkl", "rb") as f:
        model = pickle.load(f)
except FileNotFoundError:
    st.error("Model or vectorizer file not found. Please make sure the .pkl files are in the correct directory.")
    st.stop()


# ---------- Preprocessing function ----------
def clean_text(text):
    """Cleans the input text by converting to lowercase and removing non-alphabetic characters."""
    text = text.lower()
    text = re.sub(r'[^a-zA-Z\s]', '', text) # Keep spaces for word separation
    text = re.sub(r'\s+', ' ', text).strip()
    return text

# ---------- File Reader Functions ----------
def read_pdf(file):
    """Reads a PDF file using PyMuPDF and extracts text from all pages."""
    # To read a file-like object, PyMuPDF needs the bytes.
    file_bytes = file.read()
    text = ""
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def read_docx(file):
    """Reads a DOCX file and extracts text from all paragraphs."""
    doc = Document(file)
    return " ".join([p.text for p in doc.paragraphs])

def read_pptx(file):
    """Reads a PPTX file and extracts text from all shapes."""
    prs = Presentation(file)
    return " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

def read_txt(file):
    """Reads a plain text file."""
    # The uploaded file object needs to be decoded.
    return file.read().decode("utf-8")

# ---------- Streamlit App UI ----------
st.set_page_config(page_title="Fake News Detector", layout="wide", initial_sidebar_state="auto")

st.title("📰 Fake News Detector")
st.markdown("---")
st.markdown("""
Welcome to the Fake News Detector! This tool uses a machine learning model 
to analyze news content and predict whether it's likely real or fake.

**How to use:**
1.  Choose your input method: **Manual Input** or **Upload File**.
2.  Provide the news title and content, or upload a document.
3.  Click the **Detect** button to see the analysis.
""")

# --- Input Method Selection ---
st.sidebar.header("Input Method")
method = st.sidebar.radio("Choose how to provide news content:", ["Manual Input", "Upload File"])

text_data = ""

# --- Main Panel for Input ---
if method == "Manual Input":
    st.header("✍️ Manual News Entry")
    title = st.text_input("News Title", placeholder="e.g., NASA Confirms Discovery of a New Earth-Like Planet")
    text = st.text_area("News Content", height=250, placeholder="Paste the full news article content here...")
    if title and text:
        text_data = title + " " + text

else: # File Upload
    st.header("📄 Upload a Document")
    uploaded = st.file_uploader(
        "Upload a PDF, DOCX, PPTX, or TXT file",
        type=["pdf", "docx", "pptx", "txt"]
    )
    if uploaded:
        ext = uploaded.name.split(".")[-1].lower()
        try:
            with st.spinner(f"Reading {ext.upper()} file..."):
                if ext == "pdf":
                    text_data = read_pdf(uploaded)
                elif ext == "docx":
                    text_data = read_docx(uploaded)
                elif ext == "pptx":
                    text_data = read_pptx(uploaded)
                elif ext == "txt":
                    text_data = read_txt(uploaded)
            st.success("File read successfully!")
        except Exception as e:
            st.error(f"Error reading file: {e}")

# --- Analysis and Prediction ---
st.markdown("---")
if st.button("Detect", key="detect_button"):
    if not text_data or not text_data.strip():
        st.warning("Please provide some news content to analyze.")
    else:
        with st.spinner("Analyzing content..."):
            # 1. Clean the text
            cleaned_text = clean_text(text_data)
            
            # 2. Vectorize the text
            vectorized_text = vectorizer.transform([cleaned_text])
            
            # 3. Predict using the model
            prediction = model.predict(vectorized_text)

        st.header("Analysis Result")
        # Assuming 1 is 'REAL' and 0 is 'FAKE'
        if prediction[0] == 1:
            st.success("✅ This news appears to be REAL.")
        else:
            st.error("🚨 This news appears to be FAKE.")

        with st.expander("Show Analyzed Text"):
            st.text_area("Original Text", text_data, height=200)
